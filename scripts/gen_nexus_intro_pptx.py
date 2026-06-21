"""
用 python-pptx 生成 Nexus 项目介绍 PPT
- 输入:nexus/CLAUDE.md + nexus/SPEC.md 摘要
- 输出:nexus-intro.pptx
- 真实 .pptx,PowerPoint 可直接打开编辑
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 主题色 (Nexus 品牌色:深紫 + 橙)
COLOR_PRIMARY = RGBColor(0x4A, 0x2E, 0x7A)    # 深紫
COLOR_ACCENT = RGBColor(0xFF, 0x7A, 0x29)     # 橙
COLOR_TEXT = RGBColor(0x1F, 0x1F, 0x1F)
COLOR_MUTED = RGBColor(0x66, 0x66, 0x66)

# 16:9 宽屏
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_title_bar(slide, title, subtitle=None):
    """顶部标题条"""
    from pptx.shapes.connector import Connector
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=COLOR_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, size=18):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"• {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(8)
    return box


def add_accent_line(slide, x, y, w):
    """橙色彩条"""
    line = slide.shapes.add_shape(1, x, y, w, Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()


# ===== 第 1 页:封面 =====
s1 = prs.slides.add_slide(BLANK)
bg = s1.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_PRIMARY
bg.line.fill.background()

add_text(s1, "NEXUS", Inches(0.5), Inches(2.0), Inches(12), Inches(1.5),
         size=88, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s1, "AI Gateway · 智能对话 · 会话管理 · 记忆系统", Inches(0.5), Inches(3.6), Inches(12), Inches(0.6),
         size=22, color=RGBColor(0xFF, 0xC2, 0x99))
add_text(s1, "夜小白科技有限公司", Inches(0.5), Inches(5.0), Inches(12), Inches(0.5),
         size=18, color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(s1, "项目介绍 · 2026.06", Inches(0.5), Inches(6.3), Inches(12), Inches(0.4),
         size=14, color=RGBColor(0xAA, 0xAA, 0xAA))

# ===== 第 2 页:项目概述 =====
s2 = prs.slides.add_slide(BLANK)
add_title_bar(s2, "项目概述", "AI Gateway · 三进程架构")
add_accent_line(s2, Inches(0.5), Inches(1.2), Inches(2.0))
add_bullets(s2, [
    "名称:Nexus(夜小白科技有限公司开发)",
    "用途:AI Gateway,统一管理智能对话 / 会话 / 记忆 / MCP / 微信通道",
    "技术栈:React 19 + FastAPI + DeepAgents + WebSocket + SQLite + Electron",
    "三进程架构:Python 后端(30000)· React 前端(30077)· Electron 桌面端",
    "强制规范:.venv 虚拟环境 · ruff lint · pytest 测试 · 单文件 ≤800 行",
], Inches(0.7), Inches(1.6), Inches(12), Inches(4))

# ===== 第 3 页:技术架构 =====
s3 = prs.slides.add_slide(BLANK)
add_title_bar(s3, "技术架构", "后端 · 前端 · 桌面")
add_accent_line(s3, Inches(0.5), Inches(1.2), Inches(2.0))

# 三列卡片
def card(slide, x, y, w, h, title, items, color):
    box = slide.shapes.add_shape(1, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xFA)
    box.line.color.rgb = color
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = color
    for it in items:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = f"  {it}"
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_TEXT

card(s3, Inches(0.5), Inches(1.6), Inches(4.0), Inches(5.2),
     "后端 · Backend", [
         "FastAPI (端口 30000)",
         "main.py 入口",
         "agent.py DeepAgents 封装",
         "db.py SQLite + 自动迁移",
         "memory.py 记忆系统",
         "mcp.py MCP 插件",
         "channels/ 渠道层",
         "plugins/ 插件层",
     ], COLOR_PRIMARY)

card(s3, Inches(4.7), Inches(1.6), Inches(4.0), Inches(5.2),
     "前端 · Frontend", [
         "Vite + React 19 (30077)",
         "src/components/",
         "src/hooks/",
         "Zustand 状态管理",
         "WebSocket 实时",
         "E2E 测试",
         "TypeScript",
         "可独立部署",
     ], COLOR_ACCENT)

card(s3, Inches(8.9), Inches(1.6), Inches(4.0), Inches(5.2),
     "桌面 · Desktop", [
         "Electron (macOS DMG)",
         "src/main.ts 主进程",
         "src/backend.ts 后端",
         "src/preload.ts 预加载",
         "electron-builder 打包",
         "本地可执行",
         "内嵌 Python 后端",
         "可独立运行",
     ], COLOR_PRIMARY)

# ===== 第 4 页:CLI 命令 =====
s4 = prs.slides.add_slide(BLANK)
add_title_bar(s4, "CLI 命令", "服务管理 · 开发 · 测试")
add_accent_line(s4, Inches(0.5), Inches(1.2), Inches(2.0))

add_text(s4, "服务管理", Inches(0.7), Inches(1.6), Inches(6), Inches(0.4),
         size=20, bold=True, color=COLOR_PRIMARY)
add_text(s4,
         "nexus install      # 首次安装\n"
         "nexus start        # 启动\n"
         "nexus stop         # 停止\n"
         "nexus restart      # 重启\n"
         "nexus status       # 状态\n"
         "nexus logs         # 日志\n"
         "nexus doctor       # 健康检查\n"
         "nexus uninstall    # 卸载",
         Inches(0.7), Inches(2.0), Inches(6), Inches(3),
         size=14, color=COLOR_TEXT)

add_text(s4, "开发与测试", Inches(7.0), Inches(1.6), Inches(6), Inches(0.4),
         size=20, bold=True, color=COLOR_ACCENT)
add_text(s4,
         "source .venv/bin/activate\n"
         "pytest tests/                 # 全量\n"
         "pytest tests/test_xxx.py      # 单测\n"
         "ruff check nexus/             # lint\n"
         "ruff format nexus/            # 格式化\n\n"
         "cd frontend && npm install\n"
         "npm run dev|build|lint|test:e2e\n\n"
         "cd desktop && npm install\n"
         "npm run dev|test|pack",
         Inches(7.0), Inches(2.0), Inches(6), Inches(4.5),
         size=13, color=COLOR_TEXT)

# ===== 第 5 页:记忆与插件系统 =====
s5 = prs.slides.add_slide(BLANK)
add_title_bar(s5, "记忆系统 & 插件", "BM25 · SQLite · MCP 扩展")
add_accent_line(s5, Inches(0.5), Inches(1.2), Inches(2.0))

add_text(s5, "记忆系统(Memory)", Inches(0.7), Inches(1.6), Inches(12), Inches(0.5),
         size=22, bold=True, color=COLOR_PRIMARY)
add_bullets(s5, [
    "BM25 全文检索(独立于 OpenClaw 体系)",
    "上下文压缩:85% 阈值自动触发",
    "会话级 / 全局两级记忆",
    "WebSocket 实时记忆同步",
], Inches(0.7), Inches(2.2), Inches(12), Inches(2.0))

add_text(s5, "插件与渠道(Plugins & Channels)", Inches(0.7), Inches(4.4), Inches(12), Inches(0.5),
         size=22, bold=True, color=COLOR_ACCENT)
add_bullets(s5, [
    "MCP 插件:动态加载工具,扩展 agent 能力",
    "微信通道:二维码登录 + 自动会话创建",
    "DingTalk 通道:钉钉机器人推送",
    "插件白名单机制,安全可控",
], Inches(0.7), Inches(5.0), Inches(12), Inches(2.0))

# ===== 第 6 页:路线与目标 =====
s6 = prs.slides.add_slide(BLANK)
add_title_bar(s6, "近期目标", "AI 架构师方向 · 持续迭代")
add_accent_line(s6, Inches(0.5), Inches(1.2), Inches(2.0))

add_bullets(s6, [
    "PPT 文档自动化(ppt-master 集成已落地阿奇代码层)",
    "微调栈持续学习:多 LoRA 并发训练",
    "Agent 工程化:工具链 + 训练闭环",
    "西安 AI 架构师定位:3 万月薪是市场中位",
    "技术研究方向:联合博士生 / 顾问合作",
    "Nexus + Furion 双项目并行推进",
], Inches(0.7), Inches(1.8), Inches(12), Inches(4.5))

# ===== 第 7 页:联系方式 =====
s7 = prs.slides.add_slide(BLANK)
bg = s7.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_PRIMARY
bg.line.fill.background()

add_text(s7, "Thank You", Inches(0.5), Inches(2.5), Inches(12), Inches(1.5),
         size=72, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s7, "Nexus · AI Gateway", Inches(0.5), Inches(4.0), Inches(12), Inches(0.6),
         size=24, color=RGBColor(0xFF, 0xC2, 0x99))
add_text(s7, "夜小白科技有限公司 · 西安 · 2026", Inches(0.5), Inches(5.5), Inches(12), Inches(0.4),
         size=14, color=RGBColor(0xAA, 0xAA, 0xAA))
add_text(s7, "本 PPT 由 python-pptx 1.0.2 自动生成 · 7 页", Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
         size=11, color=RGBColor(0x88, 0x88, 0x88))

# 落盘
out = Path("/Users/yxb/.openclaw/workspace/diagrams/nexus-intro.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"✔ Saved: {out} ({out.stat().st_size} bytes, {len(prs.slides)} slides)")
